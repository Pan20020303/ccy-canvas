from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "outputs" / "presentations" / "2026-06-04-ccy-canvas-client-pitch"
OUT_PATH = OUT_DIR / "ccy-canvas-client-pitch.pptx"

LOGO = ROOT / "src" / "imports" / "logo.png"
SCREEN_1 = ROOT / "src" / "imports" / "image-9.png"
SCREEN_2 = ROOT / "src" / "imports" / "image-14.png"
SCREEN_3 = ROOT / "src" / "imports" / "image-22.png"
SCREEN_4 = ROOT / "src" / "imports" / "image-1.png"


BG = RGBColor(8, 11, 18)
SURFACE = RGBColor(20, 26, 38)
SURFACE_ALT = RGBColor(28, 36, 51)
TEXT = RGBColor(245, 247, 250)
TEXT_SOFT = RGBColor(182, 191, 205)
ACCENT = RGBColor(244, 79, 18)
ACCENT_SOFT = RGBColor(255, 137, 76)
BLUE = RGBColor(74, 143, 255)
LINE = RGBColor(55, 66, 86)
SUCCESS = RGBColor(78, 201, 140)

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"


def rgb(color: RGBColor) -> tuple[int, int, int]:
    return color[0], color[1], color[2]


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, line_color=None, radius=True, transparency=0):
    shape_type = SHAPE.ROUNDED_RECTANGLE if radius else SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.fill.transparency = transparency
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if radius and hasattr(shape, "adjustments"):
        shape.adjustments[0] = 0.12
    return shape


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font_size=20,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name=FONT_CN,
    linespacing=1.2,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for raw_line in text.split("\n"):
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        p.text = raw_line
        p.alignment = align
        p.line_spacing = linespacing
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
        first = False
    return box


def add_bullets(slide, x, y, w, h, items, *, font_size=20, color=TEXT_SOFT, bullet_color=ACCENT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for item in items:
        p = frame.paragraphs[0] if first else frame.add_paragraph()
        p.text = item
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.bullet = True
        p.line_spacing = 1.25
        p.space_after = Pt(10)
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = FONT_CN
        if p.runs:
            p.runs[0].font.color.rgb = color
        first = False
    return box


def add_picture_cover(slide, img_path: Path, x, y, w, h):
    slide.shapes.add_picture(str(img_path), x, y, w, h)
    border = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    border.fill.background()
    border.line.color.rgb = LINE
    border.line.width = Pt(1.2)
    if hasattr(border, "adjustments"):
        border.adjustments[0] = 0.12
    return border


def add_title_block(slide, eyebrow, title, subtitle, *, right_note=None, title_size=22):
    add_text(slide, Inches(0.7), Inches(0.45), Inches(2.6), Inches(0.3), eyebrow, font_size=12, color=ACCENT_SOFT, bold=True, font_name=FONT_EN)
    add_text(slide, Inches(0.7), Inches(0.82), Inches(10.8), Inches(0.86), title, font_size=title_size, color=TEXT, bold=True)
    add_text(slide, Inches(0.7), Inches(1.58), Inches(10.9), Inches(0.42), subtitle, font_size=12.5, color=TEXT_SOFT)
    if right_note:
        add_text(slide, Inches(9.8), Inches(0.55), Inches(2.3), Inches(0.4), right_note, font_size=11, color=TEXT_SOFT, align=PP_ALIGN.RIGHT, font_name=FONT_EN)


def add_kpi_card(slide, x, y, w, h, title, value, desc, accent=ACCENT):
    add_rect(slide, x, y, w, h, SURFACE, line_color=LINE)
    add_rect(slide, x + Inches(0.18), y + Inches(0.2), Inches(0.08), h - Inches(0.4), accent, radius=True)
    add_text(slide, x + Inches(0.42), y + Inches(0.18), w - Inches(0.6), Inches(0.22), title, font_size=12, color=TEXT_SOFT)
    add_text(slide, x + Inches(0.42), y + Inches(0.48), w - Inches(0.6), Inches(0.42), value, font_size=18, color=TEXT, bold=True)
    add_text(slide, x + Inches(0.42), y + Inches(0.9), w - Inches(0.6), h - Inches(0.98), desc, font_size=10.5, color=TEXT_SOFT)


def add_badge(slide, x, y, text, fill=ACCENT):
    width = Inches(0.55 + 0.11 * len(text))
    add_rect(slide, x, y, width, Inches(0.28), fill, radius=True)
    add_text(slide, x + Inches(0.08), y + Inches(0.04), width - Inches(0.16), Inches(0.18), text, font_size=10.5, color=TEXT, bold=True, font_name=FONT_EN)


def add_flow_step(slide, x, y, w, h, number, title, desc):
    add_rect(slide, x, y, w, h, SURFACE, line_color=LINE)
    add_rect(slide, x + Inches(0.18), y + Inches(0.18), Inches(0.42), Inches(0.42), ACCENT, radius=True)
    add_text(slide, x + Inches(0.18), y + Inches(0.21), Inches(0.42), Inches(0.2), number, font_size=12, color=TEXT, bold=True, align=PP_ALIGN.CENTER, font_name=FONT_EN)
    add_text(slide, x + Inches(0.74), y + Inches(0.18), w - Inches(0.9), Inches(0.26), title, font_size=15, color=TEXT, bold=True)
    add_text(slide, x + Inches(0.18), y + Inches(0.72), w - Inches(0.36), h - Inches(0.9), desc, font_size=11.5, color=TEXT_SOFT)


def add_timeline(slide, x, y, labels):
    line = slide.shapes.add_connector(1, x, y, x + Inches(8.2), y)
    line.line.color.rgb = LINE
    line.line.width = Pt(2)
    step_gap = Inches(2.65)
    for idx, (title, desc) in enumerate(labels):
        cx = x + step_gap * idx
        dot = slide.shapes.add_shape(SHAPE.OVAL, cx - Inches(0.12), y - Inches(0.12), Inches(0.24), Inches(0.24))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT if idx < 2 else BLUE
        dot.line.fill.background()
        add_text(slide, cx - Inches(0.2), y + Inches(0.2), Inches(1.8), Inches(0.28), title, font_size=13, color=TEXT, bold=True)
        add_text(slide, cx - Inches(0.2), y + Inches(0.55), Inches(2.2), Inches(0.5), desc, font_size=11, color=TEXT_SOFT)


def add_footer(slide, index, total):
    add_text(slide, Inches(11.7), Inches(6.9), Inches(0.5), Inches(0.2), f"{index:02d}/{total:02d}", font_size=10.5, color=TEXT_SOFT, align=PP_ALIGN.RIGHT, font_name=FONT_EN)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    total = 10

    # 1 cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_rect(slide, Inches(0.55), Inches(0.45), Inches(12.2), Inches(6.55), SURFACE_ALT, line_color=LINE)
    slide.shapes.add_picture(str(LOGO), Inches(0.85), Inches(0.8), Inches(1.0), Inches(1.1))
    add_badge(slide, Inches(2.0), Inches(0.95), "CLIENT PITCH", fill=BLUE)
    add_text(slide, Inches(0.95), Inches(2.05), Inches(5.4), Inches(1.0), "CCY Canvas", font_size=30, color=TEXT, bold=True, font_name=FONT_EN)
    add_text(slide, Inches(0.95), Inches(2.75), Inches(5.2), Inches(1.15), "让多模型 AI 生产流程\n真正进入业务体系", font_size=24, color=TEXT, bold=True)
    add_text(slide, Inches(0.95), Inches(4.15), Inches(4.8), Inches(1.0), "面向客户/合作方的试点介绍材料\n聚焦内容生产提效、资产沉淀、团队协同与可管理落地。", font_size=13, color=TEXT_SOFT)
    add_badge(slide, Inches(0.95), Inches(5.55), "商务稳重版", fill=ACCENT)
    add_badge(slide, Inches(2.25), Inches(5.55), "适合采购/试点沟通", fill=SURFACE)
    add_picture_cover(slide, SCREEN_1, Inches(6.45), Inches(0.9), Inches(5.6), Inches(4.35))
    add_text(slide, Inches(6.8), Inches(5.55), Inches(5.0), Inches(0.6), "节点式 AI 画布，把文本、图片、视频、音频流程拉回同一工作台。", font_size=13, color=TEXT_SOFT)
    add_footer(slide, 1, total)

    # 2 pain
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "01  BUSINESS PAIN", "大多数团队不是不会用 AI，而是生产流程太碎", "工具分散、模型分散、账号分散、资产分散，导致 AI 能力很难真正变成组织能力。")
    pain_cards = [
        ("工具碎片化", "图片、文本、视频分散在不同工具里，流程靠人工搬运，协作效率低。"),
        ("成果不沉淀", "Prompt、参考图、历史版本与最终产出难复用，团队反复从零开始。"),
        ("管理不可控", "谁能用什么模型、成本花到哪一步、权限怎么分配，缺乏统一入口。"),
        ("试点难放大", "个人能跑通，不代表团队能复制；一旦扩人，流程就容易失序。"),
    ]
    for idx, (title, desc) in enumerate(pain_cards):
        add_kpi_card(slide, Inches(0.8 + 3.05 * (idx % 2) + 5.95 * (idx // 2)), Inches(2.0 + 2.0 * (idx % 2 == 1)), Inches(2.75), Inches(1.55), title, "问题并不在模型", desc, accent=ACCENT if idx < 2 else BLUE)
    add_rect(slide, Inches(8.0), Inches(1.95), Inches(4.45), Inches(4.7), SURFACE, line_color=LINE)
    add_text(slide, Inches(8.35), Inches(2.25), Inches(3.6), Inches(0.45), "客户最常见的真实困境", font_size=17, color=TEXT, bold=True)
    add_bullets(slide, Inches(8.35), Inches(2.85), Inches(3.55), Inches(2.5), [
        "创意环节可以靠 AI 提速，但无法形成标准化流程。",
        "团队成员各自使用不同模型，结果质量和成本都难统一。",
        "生成结果散落在聊天记录和个人设备里，无法沉淀为组织资产。",
        "到了采购或试点阶段，缺少可部署、可管理、可扩展的系统形态。",
    ], font_size=13)
    add_rect(slide, Inches(8.35), Inches(5.65), Inches(3.4), Inches(0.65), SURFACE_ALT, line_color=LINE)
    add_text(slide, Inches(8.55), Inches(5.85), Inches(3.0), Inches(0.2), "结论：客户要的不是单点工具，而是业务可落地的生产台。", font_size=11.5, color=SUCCESS)
    add_footer(slide, 2, total)

    # 3 solution
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "02  SOLUTION", "CCY Canvas 把创意、生成、复用、协作拉回同一个工作台", "前台做创作，后台做治理，中间用统一模型接入、资产沉淀和团队空间连接。")
    blocks = [
        ("输入与规划", "文本节点\n参考图\n场景说明", Inches(0.95)),
        ("多模型生成", "文生图\n图像编辑\n文本生成\n视频扩展", Inches(3.45)),
        ("资产沉淀", "历史资产\n本地上传\n版本复用\n结果缓存", Inches(5.95)),
        ("团队协作", "团队空间\n权限分层\n成员管理\n邀请机制", Inches(8.45)),
    ]
    for title, desc, x in blocks:
        add_rect(slide, x, Inches(2.65), Inches(2.0), Inches(1.7), SURFACE, line_color=LINE)
        add_text(slide, x + Inches(0.2), Inches(2.9), Inches(1.6), Inches(0.28), title, font_size=15, color=TEXT, bold=True)
        add_text(slide, x + Inches(0.2), Inches(3.35), Inches(1.6), Inches(0.8), desc, font_size=12, color=TEXT_SOFT, align=PP_ALIGN.CENTER)
    for x in [Inches(2.95), Inches(5.45), Inches(7.95)]:
        connector = slide.shapes.add_connector(1, x, Inches(3.5), x + Inches(0.45), Inches(3.5))
        connector.line.color.rgb = ACCENT_SOFT
        connector.line.width = Pt(2)
    add_rect(slide, Inches(0.95), Inches(5.1), Inches(11.3), Inches(1.15), SURFACE_ALT, line_color=LINE)
    add_text(slide, Inches(1.2), Inches(5.42), Inches(10.7), Inches(0.28), "一句话理解：它不是聊天框，而是把 AI 生产流程做成“可编排、可复用、可协作、可管理”的工作台。", font_size=15, color=TEXT)
    add_footer(slide, 3, total)

    # 4 value
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "03  VALUE", "从单点工具，升级为可管理的内容生产系统", "客户在意的不只是能不能生成，而是能不能更快、更稳、更省、更能复制。")
    value_cards = [
        ("效率提升", "把文本、图像、视频串成流程，减少重复切换工具和重复输入。"),
        ("资产复用", "历史资产与本地缓存让优秀结果、参考图和工作流可以二次调用。"),
        ("协作落地", "团队空间、成员机制与邀请管理，让 AI 从个人尝试走向团队协作。"),
        ("治理可控", "统一模型接入和后台配置，为权限、成本与部署管理留出明确抓手。"),
    ]
    xs = [Inches(0.85), Inches(6.6), Inches(0.85), Inches(6.6)]
    ys = [Inches(2.0), Inches(2.0), Inches(4.05), Inches(4.05)]
    accents = [ACCENT, BLUE, SUCCESS, ACCENT_SOFT]
    for i, (title, desc) in enumerate(value_cards):
        add_kpi_card(slide, xs[i], ys[i], Inches(5.1), Inches(1.65), title, "看得见的业务价值", desc, accent=accents[i])
    add_footer(slide, 4, total)

    # 5 scenario
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "04  SCENARIO", "从 Prompt 到图片/视频/文本，形成内容生产闭环", "适合做创意策划、分镜设计、视觉参考生成、内容协同和批量复用。")
    steps = [
        ("01", "输入创意", "输入场景描述、角色设定、参考图与限制条件。"),
        ("02", "多节点生成", "文本、图片、视频节点按流程串联，支持逐步迭代。"),
        ("03", "沉淀资产", "优秀结果自动留存为历史资产，后续继续调用。"),
        ("04", "团队协作", "成员在同一空间内复用结果、继续编排、统一交付。"),
    ]
    for idx, step in enumerate(steps):
        add_flow_step(slide, Inches(0.85 + 3.05 * idx), Inches(2.25), Inches(2.65), Inches(1.75), *step)
    for idx in range(3):
        connector = slide.shapes.add_connector(1, Inches(3.1 + 3.05 * idx), Inches(3.1), Inches(3.45 + 3.05 * idx), Inches(3.1))
        connector.line.color.rgb = LINE
        connector.line.width = Pt(2)
    add_rect(slide, Inches(0.85), Inches(4.55), Inches(11.55), Inches(1.5), SURFACE, line_color=LINE)
    add_bullets(slide, Inches(1.1), Inches(4.9), Inches(11.0), Inches(0.9), [
        "典型客户价值：把“创意能力”转成“可执行流程”，把“个人经验”转成“团队资产”。",
        "对采购侧的意义：试点时就能看到流程效率、资产沉淀和团队协同是否成立。",
    ], font_size=13)
    add_footer(slide, 5, total)

    # 6 highlight + screenshot
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "05  PRODUCT HIGHLIGHTS", "真实界面已经具备生产台形态", "节点编排、多模态入口、画布空间和历史结果区，已经能承载客户对试点场景的理解。", title_size=20)
    add_picture_cover(slide, SCREEN_2, Inches(0.85), Inches(2.0), Inches(7.0), Inches(4.65))
    add_rect(slide, Inches(8.15), Inches(2.0), Inches(4.2), Inches(1.2), SURFACE, line_color=LINE)
    add_text(slide, Inches(8.4), Inches(2.22), Inches(3.6), Inches(0.22), "多模态节点画布", font_size=16, color=TEXT, bold=True)
    add_text(slide, Inches(8.4), Inches(2.58), Inches(3.45), Inches(0.46), "文本、图像、视频等节点按关系组织，适合表达真实生产链路。", font_size=12, color=TEXT_SOFT)
    add_rect(slide, Inches(8.15), Inches(3.45), Inches(4.2), Inches(1.2), SURFACE, line_color=LINE)
    add_text(slide, Inches(8.4), Inches(3.67), Inches(3.6), Inches(0.22), "多模型接入", font_size=16, color=TEXT, bold=True)
    add_text(slide, Inches(8.4), Inches(4.03), Inches(3.45), Inches(0.46), "支持按场景接入不同模型，降低对单一供应商的绑定风险。", font_size=12, color=TEXT_SOFT)
    add_rect(slide, Inches(8.15), Inches(4.9), Inches(4.2), Inches(1.2), SURFACE, line_color=LINE)
    add_text(slide, Inches(8.4), Inches(5.12), Inches(3.6), Inches(0.22), "历史资产与缓存", font_size=16, color=TEXT, bold=True)
    add_text(slide, Inches(8.4), Inches(5.48), Inches(3.45), Inches(0.46), "生成结果可缓存到本地上传目录，便于持久化沉淀和二次复用。", font_size=12, color=TEXT_SOFT)
    add_footer(slide, 6, total)

    # 7 team & admin
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "06  TEAM & CONTROL", "前台做创作，后台做治理，更适合团队试点", "系统不仅解决“能不能生成”，也为“谁来用、怎么管、如何扩”提供管理接口。", title_size=20)
    add_picture_cover(slide, SCREEN_3, Inches(0.85), Inches(2.0), Inches(5.6), Inches(4.7))
    add_rect(slide, Inches(6.8), Inches(2.0), Inches(5.45), Inches(4.7), SURFACE, line_color=LINE)
    add_bullets(slide, Inches(7.15), Inches(2.35), Inches(4.7), Inches(3.5), [
        "团队空间：从个人试用切到团队协作时，不需要重建工作方式。",
        "成员与邀请：适合试点期间控制进入范围、角色归属和协作边界。",
        "模型配置：为后续成本管理、能力分层和接入扩展预留统一入口。",
        "本地存储与部署：支持结果落到本地磁盘上传目录，方便客户对数据保留有预期。",
    ], font_size=13)
    add_rect(slide, Inches(7.15), Inches(5.9), Inches(4.45), Inches(0.42), SURFACE_ALT, line_color=LINE)
    add_text(slide, Inches(7.35), Inches(6.02), Inches(4.05), Inches(0.16), "这类结构更像正式产品，而不是一次性的创意演示页。", font_size=11.5, color=SUCCESS)
    add_footer(slide, 7, total)

    # 8 deployment
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "07  DELIVERY", "支持试点，也为后续部署留下扩展空间", "对客户来说，试点阶段最重要的是能快速验证价值，后续阶段最重要的是可部署、可接管、可扩展。", title_size=20)
    left_items = [
        ("本地磁盘上传", "生成图片和上传素材可保留在本地目录，便于做持久化与资产管理。"),
        ("统一 API 与后台", "模型接入、工作区、权限和资产逻辑都能收口到统一服务。"),
        ("局域网 / 内网部署路径", "项目提供明确的部署脚本与启动方式，适合试点环境落地。"),
    ]
    for i, (title, desc) in enumerate(left_items):
        add_kpi_card(slide, Inches(0.9), Inches(2.0 + 1.55 * i), Inches(4.85), Inches(1.25), title, "交付可控", desc, accent=[ACCENT, BLUE, SUCCESS][i])
    add_rect(slide, Inches(6.15), Inches(2.0), Inches(6.05), Inches(4.55), SURFACE, line_color=LINE)
    add_text(slide, Inches(6.5), Inches(2.35), Inches(5.2), Inches(0.3), "适合客户试点的原因", font_size=18, color=TEXT, bold=True)
    add_bullets(slide, Inches(6.5), Inches(2.9), Inches(5.1), Inches(2.5), [
        "先围绕一到两条内容生产流程验证效率和质量，不需要一次性铺大。",
        "系统形态已经覆盖创作、协作、管理三层，便于客户内部汇报与评估。",
        "当试点有效时，可以自然延伸到团队规模、模型治理和资产复用场景。",
    ], font_size=13)
    add_rect(slide, Inches(6.5), Inches(5.55), Inches(5.1), Inches(0.52), SURFACE_ALT, line_color=LINE)
    add_text(slide, Inches(6.72), Inches(5.72), Inches(4.7), Inches(0.18), "关键词：可试点、可部署、可管理、可扩展。", font_size=12, color=TEXT)
    add_footer(slide, 8, total)

    # 9 pilot
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_block(slide, "08  PILOT MODEL", "建议的试点合作方式：2 到 4 周完成价值验证", "先小范围聚焦高频内容场景，再评估提效、资产复用和协作效果。", title_size=20)
    add_timeline(slide, Inches(1.0), Inches(3.0), [
        ("第 1 周", "场景梳理 + 试点目标确认"),
        ("第 2 周", "流程打磨 + 模型接入验证"),
        ("第 3 周", "团队试用 + 资产沉淀观察"),
        ("第 4 周", "效果复盘 + 采购建议输出"),
    ])
    add_rect(slide, Inches(0.95), Inches(4.3), Inches(3.6), Inches(1.7), SURFACE, line_color=LINE)
    add_text(slide, Inches(1.2), Inches(4.6), Inches(3.1), Inches(0.25), "试点目标", font_size=16, color=TEXT, bold=True)
    add_text(slide, Inches(1.2), Inches(4.98), Inches(3.0), Inches(0.7), "验证是否能把一条高频内容生产流程从个人操作升级成团队标准流程。", font_size=12, color=TEXT_SOFT)
    add_rect(slide, Inches(4.85), Inches(4.3), Inches(3.6), Inches(1.7), SURFACE, line_color=LINE)
    add_text(slide, Inches(5.1), Inches(4.6), Inches(3.1), Inches(0.25), "试点评估", font_size=16, color=TEXT, bold=True)
    add_text(slide, Inches(5.1), Inches(4.98), Inches(3.0), Inches(0.7), "看效率提升、资产可复用程度、团队协同顺畅度，以及后续部署可行性。", font_size=12, color=TEXT_SOFT)
    add_rect(slide, Inches(8.75), Inches(4.3), Inches(3.45), Inches(1.7), SURFACE, line_color=LINE)
    add_text(slide, Inches(9.0), Inches(4.6), Inches(2.95), Inches(0.25), "交付结果", font_size=16, color=TEXT, bold=True)
    add_text(slide, Inches(9.0), Inches(4.98), Inches(2.85), Inches(0.7), "形成面向采购决策的结论：继续试点、扩团队，或进入正式部署评估。", font_size=12, color=TEXT_SOFT)
    add_footer(slide, 9, total)

    # 10 CTA
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_rect(slide, Inches(0.55), Inches(0.45), Inches(12.2), Inches(6.55), SURFACE_ALT, line_color=LINE)
    slide.shapes.add_picture(str(LOGO), Inches(0.95), Inches(0.95), Inches(0.8), Inches(0.9))
    add_badge(slide, Inches(1.95), Inches(1.05), "NEXT STEP", fill=ACCENT)
    add_text(slide, Inches(0.95), Inches(1.75), Inches(7.4), Inches(0.8), "下一步建议：预约演示，启动试点", font_size=24, color=TEXT, bold=True)
    add_text(slide, Inches(0.95), Inches(2.65), Inches(7.1), Inches(0.8), "如果客户已经明确有内容生产、视觉创意、分镜策划或多角色协作需求，CCY Canvas 适合尽快进入试点验证。", font_size=14, color=TEXT_SOFT)
    add_rect(slide, Inches(0.95), Inches(3.65), Inches(5.4), Inches(1.55), SURFACE, line_color=LINE)
    add_text(slide, Inches(1.25), Inches(3.95), Inches(4.8), Inches(0.26), "建议沟通话术", font_size=16, color=TEXT, bold=True)
    add_text(slide, Inches(1.25), Inches(4.35), Inches(4.7), Inches(0.52), "先用真实业务场景跑一轮，验证流程效率、资产复用和团队协同，再进入正式采购评估。", font_size=12, color=TEXT_SOFT)
    add_picture_cover(slide, SCREEN_1, Inches(7.05), Inches(1.35), Inches(4.95), Inches(3.55))
    add_rect(slide, Inches(7.05), Inches(5.25), Inches(4.95), Inches(0.8), ACCENT, line_color=ACCENT)
    add_text(slide, Inches(7.35), Inches(5.48), Inches(4.3), Inches(0.22), "CCY Canvas | 客户试点宣讲版", font_size=15, color=TEXT, bold=True)
    add_footer(slide, 10, total)

    return prs


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    main()
